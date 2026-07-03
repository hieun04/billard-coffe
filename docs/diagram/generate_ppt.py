"""
Tao file PowerPoint trinh bay chuyen de tot nghiep
Su dung: pip install python-pptx Pillow
Chay: python generate_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Mau sac
COLORS = {
    'primary': '#1565C0',      # Xanh dam
    'secondary': '#42A5F5',   # Xanh nhat
    'accent': '#FF7043',       # Cam
    'ai': '#AB47BC',          # Tim AI
    'success': '#66BB6A',     # Xanh la
    'dark': '#212121',        # Den
    'light': '#F5F5F5',      # Xam nhat
    'white': '#FFFFFF',
}

def set_slide_background(slide, color_hex):
    """Dat mau nen cho slide"""
    color = RGBColor.from_string(color_hex[1:])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle=""):
    """Slide trang bia"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Nen xanh
    set_slide_background(slide, COLORS['primary'])
    
    # Tieu de
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Phu de
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullet_points, highlight_ai=False):
    """Slide noi dung co tieu de va danh sach"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Nen trang
    set_slide_background(slide, COLORS['white'])
    
    # Thanh tieu de
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor.from_string(COLORS['primary'][1:])
    header.line.fill.background()
    
    # Tieu de
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Noi dung
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(8.6), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Kiem tra neu la sub-point
        if point.startswith("  "):
            p.text = "• " + point.strip()
            p.font.size = Pt(18)
            p.level = 1
        else:
            p.text = "• " + point
            p.font.size = Pt(20)
            p.level = 0
        
        p.font.color.rgb = RGBColor.from_string(COLORS['dark'][1:])
        p.space_after = Pt(12)
    
    # Neu la slide AI, them icon
    if highlight_ai:
        ai_badge = slide.shapes.add_textbox(Inches(8.5), Inches(0.3), Inches(1.2), Inches(0.6))
        tf = ai_badge.text_frame
        p = tf.paragraphs[0]
        p.text = "🤖 AI"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(COLORS['ai'][1:])
    
    return slide

def add_diagram_slide(prs, title, image_path):
    """Slide chua hinh anh (sơ đồ)"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Nen trang
    set_slide_background(slide, COLORS['white'])
    
    # Thanh tieu de
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor.from_string(COLORS['primary'][1:])
    header.line.fill.background()
    
    # Tieu de
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Hinh anh
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.5), Inches(9))
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Slide 2 cot"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    set_slide_background(slide, COLORS['white'])
    
    # Thanh tieu de
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor.from_string(COLORS['primary'][1:])
    header.line.fill.background()
    
    # Tieu de chinh
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Cot trai - Tieu de
    left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.3), Inches(0.5))
    tf = left_header.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(COLORS['secondary'][1:])
    
    # Cot trai - Noi dung
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(4.3), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "✓ " + item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor.from_string(COLORS['dark'][1:])
        p.space_after = Pt(8)
    
    # Cot phai - Tieu de
    right_header = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(0.5))
    tf = right_header.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(COLORS['accent'][1:])
    
    # Cot phai - Noi dung
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.1), Inches(4.3), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "✓ " + item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor.from_string(COLORS['dark'][1:])
        p.space_after = Pt(8)
    
    return slide

def add_end_slide(prs, thank_you=True):
    """Slide ket thuc"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    set_slide_background(slide, COLORS['primary'])
    
    if thank_you:
        # Cam on
        box = slide.shapes.add_textbox(Inches(0), Inches(2.5), Inches(10), Inches(1))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = "CAM ON THAY CO!"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Hoi dap
        box2 = slide.shapes.add_textbox(Inches(0), Inches(4), Inches(10), Inches(0.8))
        tf = box2.text_frame
        p = tf.paragraphs[0]
        p.text = "Q&A"
        p.font.size = Pt(32)
        p.font.color.rgb = RGBColor(255, 200, 100)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_presentation():
    """Tao presentation day du"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ===== SLIDE 1: TRANG BIA =====
    add_title_slide(prs, 
        "HE THONG QUAN LY BILLIARD CAFE\nTICH HOP AI CHATBOT",
        "Chuyen de tot nghiep\nSinh vien: [Ho va Ten] - MSSV: [So]\nGiang vien huong dan: [Ten GVHD]")
    
    # ===== SLIDE 2: MUC LUC =====
    add_content_slide(prs, "MUC LUC", [
        "1. Gioi thieu bai toan",
        "2. Phan tich he thong",
        "3. Thiet ke he thong",
        "4. Demo tinh nang",
        "5. Ket luan & Huong phat trien"
    ])
    
    # ===== SLIDE 3: GIOI THIEU BAI TOAN =====
    add_content_slide(prs, "1. GIOI THIEU BAI TOAN", [
        "Thuc trang quan ly quan billiard hien nay:",
        "  Quy trinh quan ly thu cong, thieu he thong hoa",
        "  Khong the theo doi trang thai ban theo thoi gian thuc",
        "  Khong co he thong tu van khach hang tu dong",
        "  Khong co lich su don hang chi tiet",
        "Nghi can:",
        "  Tu dong hoa quy trinh quan ly",
        "  Tich hop AI de ho tro khach hang 24/7"
    ])
    
    # ===== SLIDE 4: DAT VAN DE =====
    add_content_slide(prs, "2. DAT VAN DE", [
        "Van de 1: Khong theo doi trang thai ban theo thoi gian thuc",
        "  Ban trong/ban ban khong duoc cap nhat kip thoi",
        "Van de 2: Khach hang khong the dat ban truoc",
        "  Mat kha nang thu hut khach hang dat cho",
        "Van de 3: Khong co he thong tu van",
        "  Nhan vien phai tra loi cac cau hoi lap di lap lai",
        "Van de 4: Khong co AI ho tro",
        "  Chatbot tro giup khach hang"
    ])
    
    # ===== SLIDE 5: MUC TIEU =====
    add_content_slide(prs, "3. MUC TIEU DU AN", [
        "Xay dung phan mem quan ly ban billiard:",
        "  Quan ly danh sach ban, trang thai ban",
        "  He thong dat ban truc tuyen",
        "  Theo doi thoi gian choi va tinh phi tu dong",
        "Tich hop AI Chatbot:",
        "  Tu van khach hang 24/7",
        "  Kiem tra ban trong tu dong",
        "  Goi y mon an, uong phu hop",
        "  Phan tich du lieu va dua ra goi y"
    ])
    
    # ===== SLIDE 6: TONG QUAN HE THONG =====
    add_content_slide(prs, "4. TONG QUAN HE THONG", [
        "Cac thanh phan chinh:",
        "  Module quan ly ban",
        "  Module dat ban & Check-in/out",
        "  Module goi mon & Thanh toan",
        "  Module AI Chatbot",
        "  Module bao cao & Thong ke",
        "Nguoi dung:",
        "  Khach hang (Member)",
        "  Nhan vien (Staff)",
        "  Quan tri vien (Admin)"
    ])
    
    # ===== SLIDE 7: USE CASE =====
    img_path = "use_case_diagram.png"
    if os.path.exists(img_path):
        add_diagram_slide(prs, "5. USE CASE DIAGRAM", img_path)
    else:
        add_content_slide(prs, "5. USE CASE DIAGRAM", [
            "Actor: Member, Staff, Admin, AI Chatbot",
            "Cac chuc nang chinh:",
            "  Quan ly tai khoan (Dang ky, Dang nhap)",
            "  Quan ly dat lich (Dat ban, Huy dat)",
            "  Quan ly goi mon (Goi mon, Thanh toan)",
            "  AI Chatbot (Tu van, Kiem tra ban trong)",
            "  Quan ly he thong (Admin)"
        ])
    
    # ===== SLIDE 8: ERD =====
    add_content_slide(prs, "6. ERD - CO SO DU LIEU", [
        "Cac bang chinh:",
        "  Users - Tai khoan nguoi dung",
        "  Customers - Thong tin khach hang",
        "  Tables - Danh sach ban",
        "  Table_Sessions - Phien choi",
        "  Reservations - Danh sach dat truoc",
        "  Orders - Hoa don",
        "  Order_Items - Chi tiet hoa don",
        "  Products - San pham/mon an",
        "  AI_Conversations - Lich su chat AI"
    ])
    
    # ===== SLIDE 9: GIAO DIEN =====
    add_content_slide(prs, "7. GIAO DIEN HE THONG", [
        "Dashboard chinh:",
        "  Hien thi danh sach ban theo thoi gian thuc",
        "  Thong ke doanh thu trong ngay",
        "  So luong khach hang dang phuc vu",
        "Giao dien nguoi dung:",
        "  De su dung, thân thiện",
        "  Responsive tren dien thoai",
        "  Tich hop AI chatbot"
    ])
    
    # ===== SLIDE 10: TINH NANG DAT BAN =====
    add_two_column_slide(prs, "8. TINH NANG DAT BAN & CHECK-IN/OUT",
        "Dat ban truc tuyen",
        ["Chon ban theo vi tri", "Chon ngay va gio", "Nhap thong tin khach hang", "Xac nhan dat ban", "Nhan thong bao nhac nho"],
        "Check-in / Check-out",
        ["Quet QR hoac chon ban", "Bat dau dem gio choi", "Tu dong tinh phi thue ban", "Ket thuc phien choi", "Xuat hoa don"])
    
    # ===== SLIDE 11: GOI MON & THANH TOAN =====
    add_two_column_slide(prs, "9. GOI MON & THANH TOAN",
        "Goi mon an uong",
        ["Xem menu theo danh muc", "Them mon vao gio hang", "Cap nhat so luong", "Goi them trong khi choi", "In bill tam"],
        "Thanh toan",
        ["Tong hop phi ban + mon", "Ap dung khuyen mai", "Thanh toan nhieu hinh thuc", "Cash, Banking, QR", "Xuat hoa don dien tu"])
    
    # ===== SLIDE 12: AI CHATBOT =====
    add_content_slide(prs, "10. TINH NANG AI CHATBOT ⭐", [
        "Tu van khach hang 24/7:",
        "  Tra loi cau hoi ve ban trong",
        "  Gio mo cua, dia chi, so dien thoai",
        "  Gia ca dich vu",
        "Kiem tra & Dat ban tu dong:",
        "  Kiem tra so ban trong hien tai",
        "  Dat ban truc tiep qua chat",
        "Goi y thong minh:",
        "  Gợi y mon an / thức uống theo sở thích",
        "  Đề xuất khuyến mãi hiện có",
        "Phan tich & Du doan:",
        "  Phan tich hanh vi khach hang",
        "  Du doan gio cao diem"])
    
    # ===== SLIDE 13: CONG NGHE SU DUNG =====
    add_two_column_slide(prs, "11. CONG NGHE SU DUNG",
        "Frontend",
        ["React.js / Next.js", "TailwindCSS", "React Query"],
        "Backend",
        ["Node.js / Express", "PostgreSQL", "REST API"])
    
    # ===== SLIDE 14: KET QUA =====
    add_content_slide(prs, "12. KET QUA DAT DUOC", [
        "Cac chuc nang da hoan thanh:",
        "  He thong quan ly ban billiard",
        "  Module dat ban truc tuyen",
        "  Check-in/out tu dong",
        "  Tinh phi thue ban theo gio",
        "  Module goi mon & Thanh toan",
        "  AI Chatbot tu van khach hang",
        "  Dashboard thong ke",
        "So lieu:",
        "  45+ use cases",
        "  15 bang du lieu",
        "  Tich hop AI chatbot"])
    
    # ===== SLIDE 15: HAN CHE & HUONG PHAT TRIEN =====
    add_two_column_slide(prs, "13. HAN CHE & HUONG PHAT TRIEN",
        "Han che",
        ["Chua co thanh toan online", "AI chua hoc duoc preferences", "Chua co module kho hang"],
        "Huong phat trien",
        ["Thanh toan MoMo, ZaloPay", "AI hoc tu khach hang", "Module kho hang", "App dien thoai"])
    
    # ===== SLIDE 16: KET THUC =====
    add_end_slide(prs, thank_you=True)
    
    # Luu file
    output_path = "Chuyen_De_Tot_Nghiep_Billiard_Cafe.pptx"
    prs.save(output_path)
    print(f"Da tao xong: {output_path}")
    print(f"So slide: {len(prs.slides)}")
    
    return output_path

if __name__ == "__main__":
    create_presentation()
